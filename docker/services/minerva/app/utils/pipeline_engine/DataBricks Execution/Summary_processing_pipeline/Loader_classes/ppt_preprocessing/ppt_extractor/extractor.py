from bs4 import BeautifulSoup
from Loader_classes.ppt_preprocessing.ppt_extractor.extract_chart import ExtractChart
from Loader_classes.ppt_preprocessing.ppt_extractor.extract_connector import (
    ExtractConnector,
)
from Loader_classes.ppt_preprocessing.ppt_extractor.extract_smartart import (
    ExtractSmartart,
)
from Loader_classes.ppt_preprocessing.ppt_extractor.extract_smartart_boundary import (
    ExtractSmartartBoundary,
)
from Loader_classes.ppt_preprocessing.ppt_extractor.extract_table import ExtractTable
from Loader_classes.ppt_preprocessing.ppt_extractor.extract_textbox import (
    ExtractTextbox,
)
from Loader_classes.ppt_preprocessing.ppt_extractor.utils import (
    SlideProcessor,
    extract_smartart_xmls_from_pptx,
)
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class Extractor:
    def __init__(self, pptx_file, app_id, document_id, temp_folder_path_preprocessing):
        self.ppt_obj = Presentation(pptx_file)
        self.xmls = extract_smartart_xmls_from_pptx(pptx_file, app_id, document_id, temp_folder_path_preprocessing)

    def extract_shapes(self):
        table_number = 0
        smartart_number = 0
        slides_shapes_info = {}
        for slide_number, slide in enumerate(self.ppt_obj.slides, start=1):
            slide_shapes_info = {}
            for shape_number, shape in enumerate(slide.shapes, start=1):
                if shape.has_text_frame:
                    extract_textbox_obj = ExtractTextbox()
                    slide_shapes_info[str(shape_number) + "-textbox"] = extract_textbox_obj.extract(shape)
                elif shape.has_table:
                    table_number += 1
                    extract_table_obj = ExtractTable()
                    slide_shapes_info[str(shape_number) + "-table"] = extract_table_obj.extract(
                        shape, table_number=table_number
                    )
                elif shape.has_chart:
                    extract_chart_obj = ExtractChart()
                    slide_shapes_info[str(shape_number) + "-chart"] = extract_chart_obj.extract(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    extract_connector_obj = ExtractConnector()
                    slide_shapes_info[str(shape_number) + "-connector"] = extract_connector_obj.extract(shape)
                elif shape.shape_type:
                    pass
                # TODO: find better way to identify smartart
                else:
                    smartart_number += 1
                    extract_smartart_boundary_obj = ExtractSmartartBoundary()
                    slide_shapes_info[str(shape_number) + "-smartart_boundary"] = extract_smartart_boundary_obj.extract(
                        shape
                    )
                    if self.xmls:
                        shape_xml = self.xmls.get(smartart_number, None)
                        soup = BeautifulSoup(shape_xml, "xml")
                        smartart_shapes = soup.find_all("dsp:sp")
                        for smartart_count, smartart in enumerate(smartart_shapes, start=1):
                            extract_smartart_obj = ExtractSmartart()
                            slide_shapes_info[
                                str(shape_number) + "-smartart-" + str(smartart_count)
                            ] = extract_smartart_obj.extract(
                                smartart,
                                smartart_number=smartart_count,
                                smartart_x=slide_shapes_info[str(shape_number) + "-smartart_boundary"]["x"],
                                smartart_y=slide_shapes_info[str(shape_number) + "-smartart_boundary"]["y"],
                            )
            slides_shapes_info[slide_number] = slide_shapes_info
        return slides_shapes_info

    def process_slides(self, slides_shapes_info):
        processed_slides_shapes_info = {}
        title_slide_text = None
        slide_prefix = {}
        exclusion_words = ["agenda", "thank", "appendix"]
        for slide_number, slide_shapes_info in slides_shapes_info.items():
            slide_text_boxes_info = [
                shape_info for shape_type, shape_info in slide_shapes_info.items() if shape_info.get("text")
            ]
            slide_text_boxes_text = "\n".join(
                shape.get("text", "")
                for shape in slide_text_boxes_info
                if not shape.get("text", "").lower().startswith("proprietary")
            )

            if slide_number == 1 and len(slide_text_boxes_info) <= 5:
                print("Title Slide Detected!!")
                slide_shapes_info = {}
                title_slide_text = slide_text_boxes_text

            elif any(
                sentence.strip().lower().startswith(exclusion_word)
                for sentence in slide_text_boxes_text.split("\n")
                for exclusion_word in exclusion_words
            ):
                print("Slide Containing Exclusion words Detected!!")
                slide_shapes_info = {}

            elif (
                len(slide_text_boxes_info) <= 3
                and len(
                    [word for sentence in slide_text_boxes_text.split("\n") for word in sentence.split(" ") if word]
                )
                <= 10
            ):
                slide_prefix[slide_number + 1] = slide_text_boxes_text
                slide_shapes_info = {}

            if slide_shapes_info:
                slide_processor_obj = SlideProcessor(list(slide_shapes_info.values()))
                processed_slides_shapes_info[slide_number] = slide_processor_obj.run_all()
            else:
                processed_slides_shapes_info[slide_number] = slide_shapes_info

        return title_slide_text, slide_prefix, processed_slides_shapes_info
